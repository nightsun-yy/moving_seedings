from pathlib import Path
import sys


ROOT = Path(__file__).resolve().parents[2]
HERE = Path(__file__).resolve().parent
ASSETS = ROOT / "assets" / "ppt"
CACHE = HERE / ".cache"
OUT = HERE / "云融孪生_答辩逻辑版.pptx"
PPTX_DEPS = ROOT / ".pptx_deps"
SKILL_DIR = Path(r"C:\Users\morningstar\.codex\skills\python-pptx-builder")

sys.path.insert(0, str(PPTX_DEPS))
sys.path.insert(0, str(SKILL_DIR))

from PIL import Image, ImageOps
from ppt_maker import PresentationBuilder, Theme
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


FONT_CN = "Microsoft YaHei"
FONT_EN = "Segoe UI"
RESAMPLE = getattr(Image, "Resampling", Image).LANCZOS


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


builder = PresentationBuilder(theme=Theme.TECH)
prs = builder.prs

COLORS = {
    "blue": builder.colors["primary"],
    "blue_dark": builder.colors["secondary"],
    "cyan": builder.colors["accent"],
    "bg_dark": builder.colors["dark"],
    "bg_light": rgb("F1F7FB"),
    "white": builder.colors["text"],
    "text_dark": builder.colors["text_dark"],
    "panel_dark": rgb("1B2740"),
    "panel_mid": rgb("24364F"),
    "panel_light": rgb("EEF5FB"),
    "green": rgb("66F59A"),
    "green_dark": rgb("153B35"),
    "orange": rgb("F5AE4A"),
    "red": rgb("FF6B6B"),
    "gold": rgb("E8C46A"),
    "muted": rgb("A9BED4"),
    "muted_dark": rgb("6A788A"),
    "line": rgb("5F7793"),
}


def cover_crop(image_path: Path, ratio_tag: str, width_in: float, height_in: float) -> Path:
    CACHE.mkdir(parents=True, exist_ok=True)
    target_ratio = width_in / height_in
    out = CACHE / f"{image_path.stem}_{ratio_tag}.jpg"
    if out.exists():
        return out

    with Image.open(image_path).convert("RGB") as image:
        base = 1800
        target_size = (int(base * target_ratio), base)
        cropped = ImageOps.fit(image, target_size, method=RESAMPLE, centering=(0.5, 0.5))
        cropped.save(out, quality=92)
    return out


def add_shape(slide, shape_type, x, y, w, h, fill=None, line=None, transparency=0):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def add_picture_cover(slide, image_path: Path, x, y, w, h, ratio_tag):
    cached = cover_crop(image_path, ratio_tag, w, h)
    return slide.shapes.add_picture(str(cached), Inches(x), Inches(y), Inches(w), Inches(h))


def style_run(run, size, color, bold=False, font_name=FONT_CN):
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=18,
    color=None,
    bold=False,
    font_name=FONT_CN,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.04,
):
    color = color or COLORS["text_dark"]
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    style_run(run, size, color, bold=bold, font_name=font_name)
    return box


def add_pill(slide, x, y, w, h, text, fill, color, font_size=11, border=None):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill, line=border)
    add_text(
        slide,
        x,
        y + 0.02,
        w,
        h - 0.02,
        text,
        size=font_size,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_page_mark(slide, current, total):
    add_text(
        slide,
        12.2,
        0.2,
        0.8,
        0.2,
        f"{current:02d}/{total:02d}",
        size=10.5,
        color=COLORS["muted_dark"],
        bold=True,
        font_name=FONT_EN,
        align=PP_ALIGN.RIGHT,
    )


def add_speaker_tag(slide, speaker, time_text, dark=False):
    fill = COLORS["panel_dark"] if dark else COLORS["white"]
    color = COLORS["white"] if dark else COLORS["text_dark"]
    border = COLORS["cyan"] if dark else COLORS["blue"]
    add_pill(slide, 10.45, 0.34, 2.35, 0.34, f"{speaker} | {time_text}", fill, color, font_size=10.2, border=border)


def add_bullet_lines(slide, x, y, w, lines, color, bullet_color, size=13, line_gap=0.36):
    for idx, line in enumerate(lines):
        cy = y + idx * line_gap
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x, cy + 0.06, 0.1, 0.1, fill=bullet_color)
        add_text(slide, x + 0.16, cy, w - 0.16, 0.24, line, size=size, color=color)


def add_connector_bar(slide, x, y, w, h=0.03, color=None, alpha=0):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h, fill=color or COLORS["cyan"], transparency=alpha)


def add_chevron(slide, x, y, w, h, fill, text=None, text_color=None, size=14):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.CHEVRON, x, y, w, h, fill=fill)
    if text:
        add_text(
            slide,
            x + 0.06,
            y + 0.02,
            w - 0.12,
            h - 0.04,
            text,
            size=size,
            color=text_color or COLORS["white"],
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )


def add_card(slide, x, y, w, h, title, lines, fill, title_color, body_color, border=None, title_size=16):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill, line=border)
    add_text(slide, x + 0.18, y + 0.16, w - 0.36, 0.3, title, size=title_size, color=title_color, bold=True)
    add_bullet_lines(slide, x + 0.2, y + 0.62, w - 0.38, lines, body_color, COLORS["cyan"], size=12.2, line_gap=0.34)


def add_step_node(slide, x, y, w, h, step_num, title, lines, fill, title_color, body_color):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill, line=COLORS["cyan"])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x + 0.18, y + 0.18, 0.42, 0.42, fill=COLORS["cyan"])
    add_text(
        slide,
        x + 0.18,
        y + 0.21,
        0.42,
        0.34,
        str(step_num),
        size=15,
        color=COLORS["bg_dark"],
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(slide, x + 0.72, y + 0.16, w - 0.92, 0.3, title, size=15.5, color=title_color, bold=True)
    add_bullet_lines(slide, x + 0.2, y + 0.72, w - 0.4, lines, body_color, COLORS["green"], size=11.6, line_gap=0.28)


def add_value_card(slide, x, y, w, h, tag, title, desc, fill, text_color, tag_fill):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill)
    add_pill(slide, x + 0.2, y + 0.16, 0.95, 0.28, tag, tag_fill, COLORS["bg_dark"], font_size=9.6)
    add_text(slide, x + 0.2, y + 0.56, w - 0.4, 0.32, title, size=16, color=text_color, bold=True)
    add_text(slide, x + 0.2, y + 0.92, w - 0.4, 0.48, desc, size=11.8, color=text_color)


def add_media_card(slide, x, y, w, h, image_path: Path, title, source, url):
    add_picture_cover(slide, image_path, x, y, w, h, f"media_{image_path.stem}_{int(y * 100)}")
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, fill=COLORS["bg_dark"], transparency=0.38)
    play = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, x + 0.18, y + 0.25, 0.32, 0.42, fill=COLORS["orange"])
    play.rotation = 90
    add_text(slide, x + 0.62, y + 0.16, w - 0.75, 0.24, title, size=11.5, color=COLORS["white"], bold=True)
    add_text(slide, x + 0.62, y + 0.46, w - 0.75, 0.2, source, size=8.9, color=COLORS["bg_light"], font_name=FONT_EN)
    add_text(slide, x + 0.62, y + 0.73, w - 0.75, 0.18, "点击卡片跳转视频", size=8.5, color=COLORS["gold"], bold=True)
    click = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, fill=None)
    click.click_action.hyperlink.address = url


def add_title_header(slide, title, subtitle=None, dark=False):
    if dark:
        add_text(slide, 0.45, 0.3, 7.5, 0.38, title, size=24, color=COLORS["white"], bold=True)
        if subtitle:
            add_text(slide, 0.47, 0.72, 9.4, 0.24, subtitle, size=12.2, color=COLORS["muted"])
    else:
        add_text(slide, 0.45, 0.3, 7.5, 0.38, title, size=24, color=COLORS["text_dark"], bold=True)
        if subtitle:
            add_text(slide, 0.47, 0.72, 9.4, 0.24, subtitle, size=12.2, color=COLORS["muted_dark"])


def slide_cover():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_cover(slide, ASSETS / "cover_hydroponic.jpg", 0, 0, 13.333, 7.5, "cover_logic")
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=COLORS["bg_dark"], transparency=0.3)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 5.55, 7.5, fill=COLORS["bg_dark"], transparency=0.02)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, -0.5, 5.2, 2.6, 2.6, fill=COLORS["green"], transparency=0.76)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 3.8, -0.5, 2.1, 2.1, fill=COLORS["cyan"], transparency=0.82)
    add_pill(slide, 0.62, 0.56, 2.5, 0.36, "工业网络智能控制与维护答辩", COLORS["orange"], COLORS["bg_dark"], font_size=11)
    add_text(slide, 0.62, 1.45, 4.8, 1.2, "云融孪生", size=30, color=COLORS["white"], bold=True)
    add_text(slide, 0.62, 2.32, 4.65, 1.18, "数字孪生驱动的植物工厂\n网络化控制系统", size=24, color=COLORS["bg_light"], bold=True)
    add_text(
        slide,
        0.64,
        4.0,
        4.55,
        0.9,
        "把“底层控制、工业网络、三维孪生、AI 巡检”讲成一个完整闭环。",
        size=14.2,
        color=COLORS["bg_light"],
    )
    add_pill(slide, 0.64, 5.52, 0.95, 0.31, "采集", COLORS["green"], COLORS["bg_dark"], font_size=10)
    add_pill(slide, 1.72, 5.52, 0.95, 0.31, "映射", COLORS["cyan"], COLORS["bg_dark"], font_size=10)
    add_pill(slide, 2.80, 5.52, 0.95, 0.31, "仿真", COLORS["orange"], COLORS["bg_dark"], font_size=10)
    add_pill(slide, 3.88, 5.52, 0.95, 0.31, "控制", COLORS["gold"], COLORS["bg_dark"], font_size=10)
    add_text(slide, 0.64, 6.14, 3.2, 0.26, "主讲：同学 C（队长）", size=11.8, color=COLORS["bg_light"], bold=True)
    add_text(slide, 0.64, 6.42, 3.8, 0.24, "建议总时长：13-14 分钟", size=10.6, color=COLORS["muted"])
    add_text(slide, 8.65, 6.84, 4.1, 0.22, "本页配图：植物工厂实拍氛围图", size=9.2, color=COLORS["bg_light"], font_name=FONT_EN)
    add_page_mark(slide, 1, 11)


def slide_pain_points():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=COLORS["bg_light"])
    add_title_header(slide, "痛点破局｜为什么这件事值得做", "把“现场难看清、数据难打通、调试代价高”收束成一个明确问题。")
    add_speaker_tag(slide, "同学 C", "约 3 分钟")

    add_picture_cover(slide, ASSETS / "workers_tablet.jpg", 0.45, 1.2, 3.05, 5.45, "pain_left")
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.45, 1.2, 3.05, 5.45, fill=COLORS["bg_dark"], transparency=0.58)
    add_text(slide, 0.76, 1.62, 2.4, 0.42, "人工巡检仍依赖经验", size=18, color=COLORS["white"], bold=True)
    add_bullet_lines(
        slide,
        0.76,
        2.28,
        2.32,
        [
            "货架密集，内部状态不直观",
            "人员频繁进出，增加污染风险",
            "发现问题时往往已经晚了",
        ],
        COLORS["bg_light"],
        COLORS["orange"],
        size=11.4,
        line_gap=0.42,
    )

    center_x, center_y, center_w, center_h = 6.15, 3.15, 2.6, 0.92
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, center_x, center_y, center_w, center_h, fill=COLORS["panel_dark"], line=COLORS["cyan"])
    add_text(
        slide,
        center_x + 0.12,
        center_y + 0.18,
        center_w - 0.24,
        0.5,
        "植物工厂\n运维难题",
        size=18,
        color=COLORS["white"],
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    nodes = [
        (4.0, 1.34, "环境限制", ["栽培架密集", "内部状态难看清"], COLORS["blue"]),
        (8.95, 1.34, "卫生风险", ["频繁进出", "病菌交叉污染"], COLORS["green_dark"]),
        (4.0, 4.78, "控制黑盒", ["PLC 数据封闭", "故障排查困难"], COLORS["panel_mid"]),
        (8.95, 4.78, "试错昂贵", ["改造成本高", "实物试错损耗大"], COLORS["panel_mid"]),
    ]
    for x, y, title, lines, fill in nodes:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, 2.45, 1.05, fill=fill, line=COLORS["cyan"])
        add_text(slide, x + 0.16, y + 0.14, 2.12, 0.24, title, size=15, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER)
        add_text(
            slide,
            x + 0.16,
            y + 0.46,
            2.12,
            0.38,
            " / ".join(lines),
            size=10.6,
            color=COLORS["bg_light"],
            align=PP_ALIGN.CENTER,
        )

    add_connector_bar(slide, 5.9, 2.4, 0.55, 0.03, COLORS["blue"])
    add_connector_bar(slide, 8.7, 2.4, 0.55, 0.03, COLORS["green"])
    add_connector_bar(slide, 5.9, 5.06, 0.55, 0.03, COLORS["cyan"])
    add_connector_bar(slide, 8.7, 5.06, 0.55, 0.03, COLORS["orange"])
    add_connector_bar(slide, 6.41, 2.43, 0.03, 0.78, COLORS["line"])
    add_connector_bar(slide, 8.87, 2.43, 0.03, 0.78, COLORS["line"])
    add_connector_bar(slide, 6.41, 4.09, 0.03, 0.72, COLORS["line"])
    add_connector_bar(slide, 8.87, 4.09, 0.03, 0.72, COLORS["line"])

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 4.0, 6.52, 7.44, 0.48, fill=COLORS["panel_dark"])
    add_text(
        slide,
        4.18,
        6.63,
        7.06,
        0.22,
        "本质不是“设备不好用”，而是缺少一个把物理层、数据层、仿真层连接起来的统一入口。",
        size=11.6,
        color=COLORS["bg_light"],
        bold=True,
    )
    add_page_mark(slide, 2, 11)


def slide_innovation():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=COLORS["bg_dark"])
    add_title_header(slide, "核心创新｜从黑盒设备到可视孪生", "两条主线：实时状态映射 + 虚拟调试前置验证。", dark=True)
    add_speaker_tag(slide, "同学 C", "创新创意", dark=True)

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.55, 1.25, 5.85, 4.7, fill=COLORS["panel_mid"], line=COLORS["red"])
    add_text(slide, 0.82, 1.52, 5.2, 0.34, "传统做法", size=18, color=COLORS["red"], bold=True, align=PP_ALIGN.CENTER)
    add_bullet_lines(
        slide,
        0.96,
        2.08,
        4.9,
        [
            "巡检依赖人工经验，发现问题滞后",
            "PLC 状态封闭，跨平台看不见",
            "算法与控制逻辑只能在实物上反复试错",
            "故障定位慢，风险前移能力弱",
        ],
        COLORS["bg_light"],
        COLORS["red"],
        size=12.7,
        line_gap=0.58,
    )

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 6.92, 1.25, 5.85, 4.7, fill=COLORS["panel_mid"], line=COLORS["green"])
    add_text(slide, 7.2, 1.52, 5.2, 0.34, "云融孪生方案", size=18, color=COLORS["green"], bold=True, align=PP_ALIGN.CENTER)
    add_bullet_lines(
        slide,
        7.34,
        2.08,
        4.92,
        [
            "将 PLC / 驱动器状态映射为前端可视数据",
            "用数字孪生空间先跑逻辑，再落到实机",
            "把巡检、识别、抓取、出库串成闭环",
            "让排故、展示、扩展都拥有统一界面",
        ],
        COLORS["bg_light"],
        COLORS["green"],
        size=12.7,
        line_gap=0.58,
    )

    add_chevron(slide, 5.95, 3.0, 1.55, 0.88, COLORS["orange"], "升级", COLORS["bg_dark"], 17)

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.62, 6.18, 4.05, 0.72, fill=COLORS["blue"])
    add_text(slide, 0.84, 6.34, 3.6, 0.28, "创新 1｜工业网络状态映射", size=14.5, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 4.88, 6.18, 3.58, 0.72, fill=COLORS["cyan"])
    add_text(slide, 5.06, 6.34, 3.18, 0.28, "PLC → 网关 → 前端", size=14.3, color=COLORS["bg_dark"], bold=True, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 8.67, 6.18, 4.05, 0.72, fill=COLORS["green"])
    add_text(slide, 8.9, 6.34, 3.56, 0.28, "创新 2｜虚拟调试前置验证", size=14.5, color=COLORS["bg_dark"], bold=True, align=PP_ALIGN.CENTER)
    add_page_mark(slide, 3, 11)


def slide_architecture():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=COLORS["bg_light"])
    add_title_header(slide, "底层架构｜工业级网络与执行层", "围绕“稳定互联”讲清硬件、控制、网络三层关系。")
    add_speaker_tag(slide, "同学 A", "硬件担当")

    add_picture_cover(slide, ASSETS / "vertical_farm.jpg", 8.35, 1.2, 4.45, 5.65, "arch_photo")
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 8.35, 1.2, 4.45, 5.65, fill=COLORS["bg_dark"], transparency=0.58)
    add_text(slide, 8.64, 1.56, 3.8, 0.34, "物理执行场景", size=18, color=COLORS["white"], bold=True)
    add_bullet_lines(
        slide,
        8.64,
        2.18,
        3.62,
        [
            "工业铝型材与丝杠机构",
            "伺服驱动稳定执行",
            "为孪生系统提供真实动作底座",
        ],
        COLORS["bg_light"],
        COLORS["gold"],
        size=11.3,
        line_gap=0.42,
    )

    stack = [
        ("上位层", "HMI / 监控 / 调度", COLORS["blue"]),
        ("网络层", "Modbus TCP / 工业交换机 / 网关", COLORS["cyan"]),
        ("控制层", "PLC / 驱动器 / 24V 控制回路", COLORS["green"]),
        ("执行层", "丝杠 / 机械臂 / 传感单元", COLORS["orange"]),
    ]
    y = 1.45
    for idx, (title, desc, fill) in enumerate(stack):
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.75, y, 6.65, 1.0, fill=fill)
        add_text(slide, 1.02, y + 0.17, 1.25, 0.28, title, size=17, color=COLORS["bg_dark"], bold=True)
        add_text(slide, 2.45, y + 0.2, 4.6, 0.25, desc, size=12.3, color=COLORS["bg_dark"], bold=True)
        if idx < len(stack) - 1:
            add_chevron(slide, 3.18, y + 1.03, 1.75, 0.32, COLORS["panel_dark"], "稳定互联", COLORS["white"], 10.5)
        y += 1.35

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.75, 6.42, 6.65, 0.42, fill=COLORS["panel_dark"])
    add_text(
        slide,
        0.96,
        6.52,
        6.15,
        0.22,
        "一句话表达：标准协议把 PLC、驱动器和上位机稳定串起来，孪生界面才有真实基础。",
        size=11.2,
        color=COLORS["bg_light"],
    )
    add_page_mark(slide, 4, 11)


def slide_gateway():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=COLORS["bg_light"])
    add_title_header(slide, "建模与网关｜把机械结构和数据流打通", "这一页讲“3D 建模轻量化”和“工业协议解析”两件事。")
    add_speaker_tag(slide, "同学 B", "链路担当")

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.58, 1.22, 5.98, 2.2, fill=COLORS["white"], line=COLORS["blue"])
    add_text(slide, 0.84, 1.46, 1.9, 0.3, "建模链路", size=17, color=COLORS["blue"], bold=True)
    add_chevron(slide, 0.94, 2.1, 1.55, 0.62, COLORS["blue"], "NX 建模", size=14)
    add_chevron(slide, 2.54, 2.1, 1.86, 0.62, COLORS["cyan"], "轻量化导出", COLORS["bg_dark"], 14)
    add_chevron(slide, 4.48, 2.1, 1.62, 0.62, COLORS["green"], "前端资源", COLORS["bg_dark"], 14)
    add_text(slide, 0.96, 2.88, 5.16, 0.28, "目标：模型可看、可转、可流畅加载，不做“重模型拖死前端”。", size=11.5, color=COLORS["muted_dark"])

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.58, 3.76, 8.28, 2.42, fill=COLORS["panel_dark"], line=COLORS["cyan"])
    add_text(slide, 0.84, 4.0, 2.6, 0.3, "通信网关与数据拓扑", size=17.5, color=COLORS["white"], bold=True)

    boxes = [
        (0.92, 4.66, 1.48, 0.68, "PLC 寄存器", COLORS["blue"]),
        (2.6, 4.66, 1.65, 0.68, "Python 网关", COLORS["cyan"]),
        (4.52, 4.66, 1.78, 0.68, "解析 / 转换", COLORS["orange"]),
        (6.58, 4.66, 1.85, 0.68, "前端数据流", COLORS["green"]),
    ]
    for x, y, w, h, label, fill in boxes:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill)
        add_text(slide, x, y + 0.18, w, 0.26, label, size=13.2, color=COLORS["bg_dark"], bold=True, align=PP_ALIGN.CENTER)

    add_chevron(slide, 2.18, 4.84, 0.34, 0.33, COLORS["white"], "→", COLORS["bg_dark"], 14)
    add_chevron(slide, 4.13, 4.84, 0.34, 0.33, COLORS["white"], "→", COLORS["bg_dark"], 14)
    add_chevron(slide, 6.18, 4.84, 0.34, 0.33, COLORS["white"], "→", COLORS["bg_dark"], 14)
    add_text(slide, 0.94, 5.56, 7.56, 0.26, "把底层工业协议解析成前端可读的高频坐标、状态和日志数据。", size=11.8, color=COLORS["bg_light"])

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 9.15, 1.22, 3.68, 4.96, fill=COLORS["white"], line=COLORS["blue"])
    add_text(slide, 9.4, 1.48, 2.9, 0.28, "演示中的关键接口对象", size=16.2, color=COLORS["text_dark"], bold=True)
    add_pill(slide, 9.42, 2.08, 1.35, 0.32, "plc_monitor.py", COLORS["panel_dark"], COLORS["white"], font_size=9.4)
    add_pill(slide, 10.86, 2.08, 1.55, 0.32, "Modbus TCP", COLORS["blue"], COLORS["white"], font_size=9.4)
    add_pill(slide, 9.42, 2.52, 1.75, 0.32, "React + Three.js", COLORS["green"], COLORS["bg_dark"], font_size=9.4)
    add_pill(slide, 11.3, 2.52, 1.12, 0.32, "MES", COLORS["orange"], COLORS["bg_dark"], font_size=9.4)
    add_bullet_lines(
        slide,
        9.4,
        3.12,
        3.0,
        [
            "链路层负责把“能读到”变成“能看懂”",
            "数据层负责把“能看懂”变成“能调用”",
            "为后续 AI 巡检和联动控制留出口",
        ],
        COLORS["muted_dark"],
        COLORS["cyan"],
        size=11.2,
        line_gap=0.42,
    )
    add_text(slide, 9.4, 5.56, 3.0, 0.26, "这一页要点：桥接机械模型和通信数据。", size=10.8, color=COLORS["muted_dark"])
    add_page_mark(slide, 5, 11)


def slide_dashboard():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=COLORS["bg_dark"])
    add_title_header(slide, "孪生界面｜把状态、坐标和动作做成运维驾驶舱", "核心不是炫酷，而是让评委一眼看见“系统在工作”。", dark=True)
    add_speaker_tag(slide, "同学 C", "集成担当", dark=True)

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.62, 1.22, 12.08, 5.92, fill=COLORS["panel_dark"], line=COLORS["cyan"])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.92, 1.52, 11.48, 0.52, fill=COLORS["panel_mid"])
    add_text(slide, 1.16, 1.66, 3.9, 0.22, "智慧农业 3D 数字孪生控制台", size=15.8, color=COLORS["white"], bold=True)
    add_pill(slide, 8.28, 1.64, 0.9, 0.22, "AUTO", COLORS["green"], COLORS["bg_dark"], font_size=8.8)
    add_pill(slide, 9.32, 1.64, 1.1, 0.22, "PLC Link", COLORS["cyan"], COLORS["bg_dark"], font_size=8.8)
    add_pill(slide, 10.58, 1.64, 0.78, 0.22, "MES", COLORS["orange"], COLORS["bg_dark"], font_size=8.8)
    add_pill(slide, 11.48, 1.64, 0.76, 0.22, "AI", COLORS["gold"], COLORS["bg_dark"], font_size=8.8)

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 1.0, 2.26, 2.35, 3.76, fill=COLORS["panel_mid"])
    add_text(slide, 1.22, 2.5, 1.7, 0.24, "MES 控制面板", size=15, color=COLORS["white"], bold=True)
    add_pill(slide, 1.22, 3.0, 1.54, 0.33, "一键 AI 寻苗", COLORS["green"], COLORS["bg_dark"], font_size=9.8)
    add_pill(slide, 1.22, 3.44, 1.16, 0.33, "回原点", COLORS["cyan"], COLORS["bg_dark"], font_size=9.8)
    add_pill(slide, 1.22, 3.88, 1.54, 0.33, "日志联动", COLORS["orange"], COLORS["bg_dark"], font_size=9.8)
    add_bullet_lines(
        slide,
        1.22,
        4.42,
        1.82,
        ["状态切换", "动作触发", "异常告警", "人工干预"],
        COLORS["bg_light"],
        COLORS["green"],
        size=10.5,
        line_gap=0.31,
    )

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 3.62, 2.26, 5.12, 3.76, fill=rgb("102030"), line=COLORS["line"])
    add_text(slide, 5.15, 2.5, 2.0, 0.24, "3D 孪生工作区", size=15.5, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER)
    for ix in [4.15, 4.9, 5.65, 6.4, 7.15]:
        add_connector_bar(slide, ix, 3.24, 0.08, 1.72, COLORS["green"])
        add_connector_bar(slide, ix - 0.18, 3.1, 0.44, 0.05, COLORS["cyan"])
        add_connector_bar(slide, ix - 0.18, 4.96, 0.44, 0.05, COLORS["cyan"])
    add_connector_bar(slide, 4.1, 3.65, 3.4, 0.08, COLORS["muted_dark"])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 5.68, 3.2, 0.42, 1.08, fill=COLORS["orange"])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 5.98, 2.94, 1.26, 0.2, fill=COLORS["cyan"])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 7.02, 2.8, 0.2, 0.2, fill=COLORS["green"])
    add_text(slide, 4.12, 5.38, 4.1, 0.26, "扫描 / 识别 / 抓取 / 出库 动作在一个窗口内完成表达", size=10.6, color=COLORS["muted"])

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 9.0, 2.26, 3.12, 1.62, fill=COLORS["panel_mid"])
    add_text(slide, 9.22, 2.5, 2.5, 0.24, "实时状态卡", size=15, color=COLORS["white"], bold=True)
    add_pill(slide, 9.22, 2.92, 0.86, 0.26, "X 轴", COLORS["cyan"], COLORS["bg_dark"], font_size=8.5)
    add_pill(slide, 10.16, 2.92, 0.86, 0.26, "Y 轴", COLORS["green"], COLORS["bg_dark"], font_size=8.5)
    add_pill(slide, 11.1, 2.92, 0.86, 0.26, "Z 轴", COLORS["orange"], COLORS["bg_dark"], font_size=8.5)
    add_text(slide, 9.24, 3.32, 2.54, 0.22, "坐标高频刷新 + 环境数据联动", size=10.4, color=COLORS["muted"])

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 9.0, 4.14, 3.12, 1.88, fill=COLORS["panel_mid"])
    add_text(slide, 9.22, 4.38, 2.5, 0.24, "日志与识别区", size=15, color=COLORS["white"], bold=True)
    for i, line in enumerate(["开始 AI 智能寻苗巡检", "识别成熟苗", "执行抓取", "出库完成"]):
        add_text(slide, 9.28, 4.82 + i * 0.28, 2.6, 0.18, f"• {line}", size=9.8, color=COLORS["bg_light"])

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.95, 6.34, 11.25, 0.42, fill=COLORS["panel_mid"])
    add_text(
        slide,
        1.14,
        6.44,
        10.7,
        0.22,
        "状态看板 + 三轴坐标 + 轻量化 3D 模型 = 面向评委和运维人员都能看懂的“驾驶舱式表达”。",
        size=11.2,
        color=COLORS["bg_light"],
    )
    add_page_mark(slide, 6, 11)


def slide_demo():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=COLORS["bg_light"])
    add_title_header(slide, "现场演示｜从最稳到最高阶的防翻车策略", "把演示动作设计成阶梯，不把答辩成败押在单一步骤上。")
    add_speaker_tag(slide, "同学 C + A/B", "约 2 分钟")

    add_connector_bar(slide, 2.02, 1.86, 0.06, 4.92, COLORS["blue"])
    add_step_node(slide, 2.28, 1.26, 5.25, 1.1, 1, "数据采集", ["手动改变环境或设备状态", "坐标 / 状态在网页端同步变化"], COLORS["white"], COLORS["text_dark"], COLORS["muted_dark"])
    add_step_node(slide, 2.28, 2.52, 5.25, 1.1, 2, "实物控制", ["触摸屏或按钮触发动作", "体现电气控制的稳定性"], COLORS["white"], COLORS["text_dark"], COLORS["muted_dark"])
    add_step_node(slide, 2.28, 3.78, 5.25, 1.1, 3, "虚实同步", ["网页下发简单动作指令", "实机响应形成云端协同"], COLORS["white"], COLORS["text_dark"], COLORS["muted_dark"])
    add_step_node(slide, 2.28, 5.04, 5.25, 1.26, 4, "虚拟调试保底", ["若实机异常，立即切换 3D 模型独立巡检动画", "把失败点转化为“虚拟调试价值”"], COLORS["panel_dark"], COLORS["white"], COLORS["bg_light"])

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 7.82, 1.24, 4.92, 5.12, fill=COLORS["panel_dark"], line=COLORS["cyan"])
    add_text(slide, 8.1, 1.5, 4.24, 0.3, "备用视频 / 行业参考", size=17, color=COLORS["white"], bold=True)
    add_text(slide, 8.1, 1.86, 4.2, 0.24, "点击卡片可跳转公开视频，用于现场保底或增强说服力。", size=10.6, color=COLORS["muted"])
    add_media_card(
        slide,
        8.08,
        2.3,
        4.34,
        1.06,
        ASSETS / "vertical_farm.jpg",
        "全自动垂直农场案例",
        "Jungheinrich AG / YouTube",
        "https://www.youtube.com/watch?v=xq_S_5BkUCI",
    )
    add_media_card(
        slide,
        8.08,
        3.62,
        4.34,
        1.06,
        ASSETS / "workers_tablet.jpg",
        "PhenoBot 温室表型机器人",
        "Wageningen UR / YouTube",
        "https://www.youtube.com/watch?v=Se6sROZoTUs",
    )
    add_media_card(
        slide,
        8.08,
        4.94,
        4.34,
        1.06,
        ASSETS / "conference_stage.jpg",
        "Greenhouse Industry 4.0 Digital Twin",
        "SDU / YouTube",
        "https://www.youtube.com/watch?v=EyBqsKfwGCc",
    )
    add_text(slide, 8.1, 6.34, 4.18, 0.28, "建议现场另行录制一段“本队设备 + 本队网页”的无剪辑备份视频。", size=9.8, color=COLORS["bg_light"])
    add_page_mark(slide, 7, 11)


def slide_value():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=COLORS["bg_light"])
    add_title_header(slide, "应用价值｜不是展示，而是能落地的能力集合", "从实用性、经济性、可持续性三个角度收束价值表达。")
    add_speaker_tag(slide, "同学 C", "约 1.5 分钟")

    add_picture_cover(slide, ASSETS / "seedlings_greenhouse.jpg", 0.58, 1.2, 4.15, 5.7, "value_left")
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.58, 1.2, 4.15, 5.7, fill=COLORS["bg_dark"], transparency=0.68)
    add_text(slide, 0.88, 1.58, 3.5, 0.38, "可落地场景", size=18, color=COLORS["white"], bold=True)
    add_bullet_lines(
        slide,
        0.88,
        2.24,
        3.22,
        ["无菌育苗中心", "高附加值作物种植舱", "低传感器初装场景", "数字化改造起步阶段"],
        COLORS["bg_light"],
        COLORS["gold"],
        size=11.6,
        line_gap=0.43,
    )

    add_value_card(
        slide,
        5.12,
        1.3,
        7.6,
        1.4,
        "01",
        "实用性｜先把“看不见”变成“可监控”",
        "即便企业尚未部署昂贵视觉传感器，也能先用孪生底座实现数字化监控的从无到有。",
        COLORS["white"],
        COLORS["text_dark"],
        COLORS["green"],
    )
    add_value_card(
        slide,
        5.12,
        3.04,
        7.6,
        1.4,
        "02",
        "经济性｜把试错成本从现场前移到虚拟空间",
        "虚拟调试减少实物试错和频繁改造带来的时间、人力与硬件损耗。",
        COLORS["white"],
        COLORS["text_dark"],
        COLORS["cyan"],
    )
    add_value_card(
        slide,
        5.12,
        4.78,
        7.6,
        1.4,
        "03",
        "可持续｜精准控制更接近绿色低碳目标",
        "网络化智控让资源投放更精确，适合高密度、节水、连续生产的现代农业方向。",
        COLORS["white"],
        COLORS["text_dark"],
        COLORS["orange"],
    )

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 5.12, 6.44, 7.6, 0.42, fill=COLORS["panel_dark"])
    add_text(
        slide,
        5.34,
        6.54,
        7.08,
        0.22,
        "行业参考：Jungheinrich 2025 垂直农场案例提到闭环系统可节水高达 95%，说明“自动化 + 可视控制”具备真实经营价值。",
        size=10.6,
        color=COLORS["bg_light"],
    )
    add_page_mark(slide, 8, 11)


def slide_teamwork():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=COLORS["bg_dark"])
    add_title_header(slide, "团队协作｜把一个延迟问题拆成三人闭环排故", "这一页不要空讲合作，要讲一次真实的工程协同路径。", dark=True)
    add_speaker_tag(slide, "同学 C", "协作与素养", dark=True)

    add_pill(slide, 0.7, 1.2, 2.5, 0.34, "问题触发：前端坐标严重延迟", COLORS["orange"], COLORS["bg_dark"], font_size=10.4)
    steps = [
        (0.88, 2.2, "A", "检查 PLC 扫描周期", COLORS["blue"]),
        (3.28, 2.2, "B", "抓包定位网关轮询瓶颈", COLORS["cyan"]),
        (6.22, 2.2, "C", "优化前端请求机制", COLORS["green"]),
        (9.34, 2.2, "OK", "虚实同步延迟显著下降", COLORS["orange"]),
    ]
    widths = [2.1, 2.55, 2.55, 2.95]
    for idx, (x, y, tag, title, fill) in enumerate(steps):
        w = widths[idx]
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, 1.12, fill=fill)
        add_pill(slide, x + 0.16, y + 0.16, 0.54 if tag != "OK" else 0.74, 0.26, tag, COLORS["bg_dark"], COLORS["white"], font_size=9.2)
        add_text(slide, x + 0.18, y + 0.54, w - 0.36, 0.32, title, size=13.4, color=COLORS["bg_dark"], bold=True, align=PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            add_chevron(slide, x + w + 0.14, y + 0.38, 0.44, 0.3, COLORS["white"], "→", COLORS["bg_dark"], 12)

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.9, 4.16, 11.7, 1.54, fill=COLORS["panel_mid"], line=COLORS["cyan"])
    add_text(
        slide,
        1.2,
        4.44,
        11.1,
        0.3,
        "这段故事要传递的不是“我们很努力”，而是：硬件、链路、界面三个人都知道自己在系统中的职责边界，并且能在同一个问题上快速闭环。",
        size=15.2,
        color=COLORS["white"],
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    add_card(
        slide,
        1.05,
        6.0,
        3.55,
        0.92,
        "A｜硬件维度",
        ["确认控制周期与底层执行稳定"],
        COLORS["panel_mid"],
        COLORS["cyan"],
        COLORS["bg_light"],
        border=COLORS["blue"],
        title_size=14.2,
    )
    add_card(
        slide,
        4.9,
        6.0,
        3.55,
        0.92,
        "B｜链路维度",
        ["抓包、定位、拆解通信瓶颈"],
        COLORS["panel_mid"],
        COLORS["green"],
        COLORS["bg_light"],
        border=COLORS["cyan"],
        title_size=14.2,
    )
    add_card(
        slide,
        8.75,
        6.0,
        3.55,
        0.92,
        "C｜界面维度",
        ["优化请求与呈现，缩短体验延迟"],
        COLORS["panel_mid"],
        COLORS["orange"],
        COLORS["bg_light"],
        border=COLORS["orange"],
        title_size=14.2,
    )
    add_page_mark(slide, 9, 11)


def slide_quality():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=COLORS["bg_light"])
    add_title_header(slide, "工程素养｜把风险前移，把细节做实", "安全和规范不是附属项，而是系统可信度的一部分。")
    add_speaker_tag(slide, "同学 C", "职业素养")

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.72, 1.32, 5.9, 4.95, fill=COLORS["white"], line=COLORS["blue"])
    add_text(slide, 1.0, 1.6, 2.2, 0.3, "电气安全", size=18, color=COLORS["blue"], bold=True)
    add_bullet_lines(
        slide,
        1.02,
        2.18,
        5.0,
        [
            "强弱电分离，减少干扰与误触风险",
            "控制回路使用 24V 安全电压",
            "PLC 底层设置双重软限位",
            "先保护设备与作物，再追求动作速度",
        ],
        COLORS["text_dark"],
        COLORS["green"],
        size=12.2,
        line_gap=0.52,
    )
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 1.02, 4.98, 4.96, 0.68, fill=COLORS["panel_dark"])
    add_text(slide, 1.24, 5.16, 4.52, 0.24, "敬畏硬件 = 在底层先把最坏情况考虑进去", size=12.2, color=COLORS["bg_light"], bold=True)

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 6.92, 1.32, 5.68, 4.95, fill=COLORS["panel_dark"], line=COLORS["cyan"])
    add_text(slide, 7.2, 1.6, 2.2, 0.3, "代码规范", size=18, color=COLORS["white"], bold=True)
    add_bullet_lines(
        slide,
        7.22,
        2.18,
        4.82,
        [
            "通信代码做异常处理，不把故障静默吞掉",
            "前端模型极致轻量化，保障演示流畅",
            "日志可回看，现场问题可追溯",
            "预设保底策略，不让演示被单点故障击穿",
        ],
        COLORS["bg_light"],
        COLORS["cyan"],
        size=12.2,
        line_gap=0.52,
    )
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 7.22, 4.98, 4.96, 0.68, fill=COLORS["panel_mid"])
    add_text(slide, 7.44, 5.16, 4.52, 0.24, "工匠精神 = 每一行代码、每一个模型面数都算数", size=12.2, color=COLORS["white"], bold=True)

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.72, 6.54, 11.88, 0.42, fill=COLORS["white"], line=COLORS["line"])
    add_text(
        slide,
        0.94,
        6.64,
        11.42,
        0.22,
        "这页收束金句：系统之所以可信，不只因为它能动，更因为它在异常、边界和风险面前也有秩序。",
        size=11.3,
        color=COLORS["muted_dark"],
    )
    add_page_mark(slide, 10, 11)


def slide_thanks():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_cover(slide, ASSETS / "conference_stage.jpg", 0, 0, 13.333, 7.5, "thanks_stage")
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=COLORS["bg_dark"], transparency=0.36)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 2.25, 1.38, 8.86, 4.72, fill=COLORS["panel_dark"], transparency=0.18, line=COLORS["cyan"])
    add_text(slide, 0, 2.55, 13.333, 0.72, "感谢聆听", size=30, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 0, 3.45, 13.333, 0.42, "欢迎各位评委专家批评指正", size=18, color=COLORS["bg_light"], align=PP_ALIGN.CENTER)
    add_pill(slide, 4.02, 4.5, 1.45, 0.32, "A 硬件执行", COLORS["green"], COLORS["bg_dark"], font_size=9.6)
    add_pill(slide, 5.94, 4.5, 1.45, 0.32, "B 链路网关", COLORS["cyan"], COLORS["bg_dark"], font_size=9.6)
    add_pill(slide, 7.86, 4.5, 1.45, 0.32, "C 孪生集成", COLORS["orange"], COLORS["bg_dark"], font_size=9.6)
    add_text(slide, 0, 6.55, 13.333, 0.24, "附：上一页视频卡片可作为现场保底演示入口", size=10.5, color=COLORS["bg_light"], align=PP_ALIGN.CENTER)
    add_page_mark(slide, 11, 11)


def build():
    slide_cover()
    slide_pain_points()
    slide_innovation()
    slide_architecture()
    slide_gateway()
    slide_dashboard()
    slide_demo()
    slide_value()
    slide_teamwork()
    slide_quality()
    slide_thanks()
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    build()
