from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / ".pptx_deps"))

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets" / "ppt"
CACHE = ASSETS / ".cache"
OUT = ROOT / "docs" / "智慧育苗系统_五页场景方案.pptx"

FONT_CN = "Microsoft YaHei"
FONT_EN = "Segoe UI"


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


COLORS = {
    "dark": rgb("173B2F"),
    "dark_2": rgb("214631"),
    "light": rgb("F5F2EA"),
    "light_2": rgb("FBF8F1"),
    "sage": rgb("D6E2CA"),
    "olive": rgb("7F9968"),
    "accent": rgb("E7862D"),
    "red": rgb("C84E41"),
    "gold": rgb("CFAE67"),
    "text": rgb("1E2D23"),
    "muted": rgb("5F6E66"),
    "white": rgb("FFFDF8"),
}


RESAMPLE = getattr(Image, "Resampling", Image).LANCZOS


def cover_crop(image_path: Path, ratio_tag: str, width_in: float, height_in: float) -> Path:
    CACHE.mkdir(parents=True, exist_ok=True)
    target_ratio = width_in / height_in
    out = CACHE / f"{image_path.stem}_{ratio_tag}.jpg"
    if out.exists():
        return out

    with Image.open(image_path).convert("RGB") as image:
        base = 1800
        target_size = (int(base * target_ratio), base)
        cropped = ImageOps.fit(image, target_size, method=RESAMPLE, centering=(0.5, 0.5))
        cropped.save(out, quality=92)
    return out


def add_shape(slide, shape_type, x, y, w, h, fill=None, line=None, radius=False, transparency=0):
    if radius:
        shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def add_picture_cover(slide, image_path: Path, x, y, w, h, ratio_tag):
    cached = cover_crop(image_path, ratio_tag, w, h)
    return slide.shapes.add_picture(str(cached), Inches(x), Inches(y), Inches(w), Inches(h))


def set_font(run, size, color, bold=False, name=FONT_CN):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=18,
    color=None,
    bold=False,
    font_name=FONT_CN,
    align=PP_ALIGN.LEFT,
    margin=0.0,
    valign=MSO_ANCHOR.TOP,
):
    color = color or COLORS["text"]
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    set_font(run, size, color, bold=bold, name=font_name)
    return box


def add_rich_lines(slide, x, y, w, lines, color=None, size=12.5, line_gap=0.33, bullet_color=None):
    color = color or COLORS["muted"]
    bullet_color = bullet_color or COLORS["accent"]
    for idx, line in enumerate(lines):
        cy = y + idx * line_gap
        bullet = add_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.OVAL,
            x,
            cy + 0.07,
            0.11,
            0.11,
            fill=bullet_color,
        )
        bullet.line.fill.background()
        add_text(slide, x + 0.18, cy, w - 0.18, 0.24, line, size=size, color=color)


def add_pill(slide, x, y, w, h, text, fill, color, font_size=10.5):
    pill = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill)
    add_text(slide, x, y + 0.03, w, h - 0.02, text, size=font_size, color=color, bold=True, align=PP_ALIGN.CENTER)
    return pill


def add_scene_row(
    slide,
    top,
    title,
    subtitle,
    bullets,
    image_path: Path,
    image_side="right",
    accent_label="典型落地",
    footer_text="",
):
    card_x, card_w, card_h = 0.35, 12.63, 2.72
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, card_x, top, card_w, card_h, fill=COLORS["white"], line=COLORS["sage"])

    img_w = 4.55
    img_h = 2.22
    img_y = top + 0.25
    if image_side == "right":
        img_x = 8.08
        text_x = 0.72
    else:
        img_x = 0.72
        text_x = 5.5

    add_picture_cover(slide, image_path, img_x, img_y, img_w, img_h, f"{image_path.stem}_{image_side}_{int(top*100)}")
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, img_x, img_y, img_w, img_h, fill=COLORS["dark"], transparency=0.78)

    add_pill(slide, text_x, top + 0.22, 1.22, 0.33, accent_label, COLORS["accent"], COLORS["white"])
    add_text(slide, text_x, top + 0.58, 6.9, 0.42, title, size=21, color=COLORS["dark"], bold=True)
    add_text(slide, text_x, top + 1.01, 6.9, 0.42, subtitle, size=12.5, color=COLORS["muted"])
    add_rich_lines(slide, text_x, top + 1.43, 6.7, bullets, size=12.3, line_gap=0.34)

    if footer_text:
        footer_y = top + 2.25
        footer_box = add_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            text_x,
            footer_y,
            6.85,
            0.26,
            fill=COLORS["sage"],
        )
        footer_box.fill.transparency = 0.1
        add_text(slide, text_x + 0.12, footer_y + 0.02, 6.56, 0.2, footer_text, size=10.5, color=COLORS["dark_2"], bold=True)


def add_page_mark(slide, num):
    add_text(slide, 12.25, 0.25, 0.6, 0.22, num, size=10.5, color=COLORS["olive"], bold=True, font_name=FONT_EN, align=PP_ALIGN.RIGHT)


def add_video_card(slide, x, y, w, h, image_path, title, source, url):
    add_picture_cover(slide, image_path, x, y, w, h, f"video_{image_path.stem}_{int(y*100)}")
    overlay = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, fill=COLORS["dark"], transparency=0.35)
    overlay.line.fill.background()

    play = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, x + 0.2, y + 0.35, 0.38, 0.5, fill=COLORS["accent"])
    play.rotation = 90
    play.line.fill.background()

    add_text(slide, x + 0.72, y + 0.2, w - 0.95, 0.28, title, size=13.5, color=COLORS["white"], bold=True)
    add_text(slide, x + 0.72, y + 0.58, w - 0.95, 0.23, source, size=10.2, color=COLORS["light"])
    add_text(slide, x + 0.72, y + 0.92, w - 0.95, 0.2, "点击图片跳转在线视频", size=9.6, color=COLORS["gold"], bold=True)

    click = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, fill=None)
    click.click_action.hyperlink.address = url
    return click


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=COLORS["light"])
    add_picture_cover(slide, ASSETS / "cover_hydroponic.jpg", 5.35, 0, 7.98, 7.5, "cover_full")
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 5.65, 7.5, fill=COLORS["dark"])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 5.35, 0, 7.98, 7.5, fill=COLORS["dark"], transparency=0.78)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, -0.5, 5.4, 2.8, 2.8, fill=COLORS["olive"], transparency=0.45)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 3.8, -0.6, 2.2, 2.2, fill=COLORS["accent"], transparency=0.6)
    add_pill(slide, 0.6, 0.58, 1.45, 0.36, "5页版方案", COLORS["accent"], COLORS["white"], font_size=11)
    add_text(slide, 0.62, 1.35, 4.5, 1.2, "智慧育苗搬运系统", size=28, color=COLORS["white"], bold=True)
    add_text(slide, 0.62, 2.28, 4.75, 1.1, "五大应用场景演示稿", size=24, color=COLORS["light"], bold=True)
    add_text(
        slide,
        0.62,
        3.58,
        4.45,
        1.2,
        "从无人巡检、AI 分拣到数字展厅，\n用一套 3D 数字孪生控制台覆盖生产、科研与展示。",
        size=15.5,
        color=COLORS["light"],
    )
    add_pill(slide, 0.62, 5.55, 1.05, 0.33, "AI巡检", COLORS["sage"], COLORS["dark_2"])
    add_pill(slide, 1.78, 5.55, 1.28, 0.33, "机械臂出库", COLORS["sage"], COLORS["dark_2"])
    add_pill(slide, 3.18, 5.55, 1.12, 0.33, "环境联动", COLORS["sage"], COLORS["dark_2"])
    add_pill(slide, 0.62, 6.02, 1.22, 0.33, "病害拦截", COLORS["gold"], COLORS["dark_2"])
    add_pill(slide, 1.96, 6.02, 1.44, 0.33, "科研采样", COLORS["gold"], COLORS["dark_2"])
    add_pill(slide, 3.52, 6.02, 1.25, 0.33, "展厅演示", COLORS["gold"], COLORS["dark_2"])
    add_text(slide, 5.88, 6.78, 6.6, 0.24, "图片来源：Pexels 实拍温室与室内农业照片", size=9.6, color=COLORS["light"], font_name=FONT_EN)
    add_page_mark(slide, "01/05")

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=COLORS["light"])
    add_text(slide, 0.45, 0.32, 4.5, 0.35, "场景 1-2｜生产端无人化", size=24, color=COLORS["dark"], bold=True)
    add_text(slide, 0.47, 0.72, 8.2, 0.25, "把“巡检困难、分级低效、病害扩散”三类痛点，收束成一套自动化作业闭环。", size=12.5, color=COLORS["muted"])
    add_scene_row(
        slide,
        1.12,
        "场景一  全自动植物工厂 / 垂直农业",
        "痛点：多层密集货架巡检困难，人工进出破坏无菌环境，夜间与非工作时段更难全覆盖。",
        [
            "机械臂携带 AI 摄像头执行夜间 S 型自动巡检，减少人员频繁进出。",
            "物理世界的幼苗状态实时映射到 3D 控制台，用绿 / 黄 / 红快速判断风险。",
            "管理人员在办公室即可总览全厂几万株植物的状态与空间分布。",
        ],
        ASSETS / "vertical_farm.jpg",
        image_side="right",
        accent_label="无人巡检",
        footer_text="关键词：多层矩阵货架｜夜间巡检｜数字孪生总览",
    )
    add_scene_row(
        slide,
        4.02,
        "场景二  规模化智能育苗基地",
        "痛点：出圃前分级高度依赖人工肉眼识别与手工搬运，效率低且质量波动大。",
        [
            "边缘 AI 自动识别成熟幼苗（状态 2），指挥 Z 轴机械爪进行精准抓取。",
            "识别到缺水或病害（状态 3 / 红色）时立即在日志中报警并执行拦截。",
            "成熟苗自动搬运至一楼出库口（0,0 坐标），形成标准化分拣出库流程。",
        ],
        ASSETS / "seedlings_greenhouse.jpg",
        image_side="left",
        accent_label="AI分拣出库",
        footer_text="关键词：寻苗｜抓取｜出库｜病害拦截",
    )
    add_page_mark(slide, "02/05")

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=COLORS["light_2"])
    add_text(slide, 0.45, 0.32, 5.2, 0.35, "场景 3-4｜科研与高价值作物", size=24, color=COLORS["dark"], bold=True)
    add_text(slide, 0.47, 0.72, 8.0, 0.25, "从长期表型采样到高价值作物保护，这一套系统不仅能“干活”，还能“保命”和“留数据”。", size=12.5, color=COLORS["muted"])
    add_scene_row(
        slide,
        1.12,
        "场景三  农业科研院所 / 表型分析平台",
        "痛点：人工测量误差大，且难以连续记录同一株植物在不同温湿光条件下的生长曲线。",
        [
            "输入特定 X、Y 坐标，驱动机械臂对单一穴位进行长期、定点图像与三维数据采集。",
            "顶部温、湿、光数据与本次抓拍图像绑定，形成可回溯的科研记录。",
            "支持 24 小时高频采样，适合育种、表型分析与环境响应研究。",
        ],
        ASSETS / "workers_tablet.jpg",
        image_side="right",
        accent_label="定点高频采集",
        footer_text="关键词：坐标采样｜EnvData 绑定｜时间戳数据集",
    )
    add_scene_row(
        slide,
        4.02,
        "场景四  高附加值经济作物种植舱",
        "痛点：中药材、航天育种等高敏作物单株价值高，一旦死苗或误碰，损失十分直接。",
        [
            "控制台提供硬件急停与原点复位，在机械臂干涉或失控时优先保护作物与设备。",
            "MES 实时监测异常植株，一旦出现红色报警可联动 PLC 执行局部滴灌或补光。",
            "适用于名贵中药材、石斛、兰科植物及其他高价值精细化种植场景。",
        ],
        ASSETS / "orchid.jpg",
        image_side="left",
        accent_label="急停与容灾",
        footer_text="关键词：急停｜原点复位｜局部滴灌｜局部补光",
    )
    add_page_mark(slide, "03/05")

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_picture_cover(slide, ASSETS / "conference_stage.jpg", 0, 0, 13.333, 7.5, "showroom_full")
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=COLORS["dark"], transparency=0.42)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.48, 0.58, 6.15, 5.95, fill=COLORS["dark"], transparency=0.16)
    add_pill(slide, 0.82, 0.88, 1.65, 0.35, "场景五｜数字展厅", COLORS["accent"], COLORS["white"], font_size=11)
    add_text(slide, 0.82, 1.46, 5.1, 0.85, "把农业系统讲成\n“看得见的科技能力”", size=26, color=COLORS["white"], bold=True)
    add_text(
        slide,
        0.82,
        2.52,
        5.4,
        1.05,
        "传统农业大屏多为 2D 图表与数字堆砌；\n而 3D 控制台能把巡检、抓取、报警和坐标数据同时“演出来”。",
        size=15.2,
        color=COLORS["light"],
    )
    add_rich_lines(
        slide,
        0.82,
        3.88,
        5.35,
        [
            "85 寸大屏投放时，动态光锥扫描与实时坐标更新能显著提升科技观感。",
            "演讲者在平板点击“一键 AI 寻苗”，玻璃房内真实机械臂可同步动作。",
            "适合投资路演、客户参观、政府汇报与企业品牌展陈。",
        ],
        color=COLORS["light"],
        bullet_color=COLORS["gold"],
        size=12.6,
        line_gap=0.37,
    )
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 7.55, 0.92, 5.2, 2.55, fill=COLORS["light"])
    add_picture_cover(slide, ASSETS / "control_room.jpg", 7.82, 1.18, 1.92, 2.0, "control_strip")
    add_text(slide, 9.98, 1.18, 2.35, 0.28, "展厅讲解动线", size=16, color=COLORS["dark"], bold=True)
    add_pill(slide, 9.98, 1.67, 2.05, 0.34, "01 大屏总览开场", COLORS["sage"], COLORS["dark_2"])
    add_pill(slide, 9.98, 2.11, 1.92, 0.34, "02 AI寻苗演示", COLORS["gold"], COLORS["dark_2"])
    add_pill(slide, 9.98, 2.55, 2.22, 0.34, "03 报警与联动收束", COLORS["accent"], COLORS["white"])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 7.55, 3.78, 5.2, 2.05, fill=COLORS["light"])
    add_text(slide, 7.9, 4.06, 4.4, 0.28, "为什么这个场景最“见效”", size=16, color=COLORS["dark"], bold=True)
    add_rich_lines(
        slide,
        7.92,
        4.48,
        4.4,
        [
            "它最适合把“种植能力”升级成“科技能力”的对外表达。",
            "既能展示系统逻辑，也能直接带动真实设备，形成虚实互动。",
            "在商务沟通场景里，3D 演示比单纯图表更容易建立记忆点。",
        ],
        color=COLORS["muted"],
        bullet_color=COLORS["accent"],
        size=12.1,
        line_gap=0.34,
    )
    add_text(slide, 0.82, 6.77, 6.4, 0.22, "背景图片：Pexels 实拍会议与屏幕场景", size=9.6, color=COLORS["light"], font_name=FONT_EN)
    add_page_mark(slide, "04/05")

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=COLORS["light"])
    add_text(slide, 0.45, 0.32, 6.2, 0.35, "一套控制台，覆盖五类业务价值", size=24, color=COLORS["dark"], bold=True)
    add_text(slide, 0.47, 0.72, 7.1, 0.25, "真正可讲的不是“炫酷 3D”，而是它把生产、科研、容灾和展示连接成了一个统一入口。", size=12.5, color=COLORS["muted"])

    left_x = 0.45
    card_w = 8.0
    y_positions = [1.2, 2.62, 4.04]
    titles = ["生产提效", "科研取数", "品牌表达"]
    labels = [
        "无人巡检｜AI 分拣出库｜病害拦截",
        "定点采样｜环境绑定｜生长曲线沉淀",
        "3D 孪生演示｜实机联动｜展厅路演",
    ]
    descriptions = [
        "解决多层货架巡检难、人工分级慢、出库不标准的问题。",
        "把图像、三维形态和温湿光数据绑定成连续科研数据集。",
        "让投资人、客户和主管部门快速看懂系统价值与技术成熟度。",
    ]
    colors = [COLORS["sage"], COLORS["gold"], COLORS["accent"]]
    text_colors = [COLORS["dark_2"], COLORS["dark_2"], COLORS["white"]]

    for idx, y in enumerate(y_positions):
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left_x, y, card_w, 1.12, fill=COLORS["white"], line=COLORS["sage"])
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left_x, y, 1.42, 1.12, fill=colors[idx])
        add_text(slide, left_x + 0.12, y + 0.28, 1.15, 0.24, titles[idx], size=15.5, color=text_colors[idx], bold=True)
        add_text(slide, left_x + 1.65, y + 0.17, 5.95, 0.24, labels[idx], size=11.5, color=COLORS["dark"], bold=True)
        add_text(slide, left_x + 1.65, y + 0.52, 5.95, 0.34, descriptions[idx], size=11.2, color=COLORS["muted"])

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.45, 5.62, 8.0, 1.18, fill=COLORS["dark"])
    add_text(slide, 0.7, 5.88, 2.1, 0.24, "建议现场演示顺序", size=15.2, color=COLORS["white"], bold=True)
    add_text(
        slide,
        2.7,
        5.84,
        5.2,
        0.55,
        "总览大屏 → AI 寻苗 → 抓取出库 → 病害告警 → 科研数据回看",
        size=13.2,
        color=COLORS["light"],
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    add_text(slide, 9.0, 0.52, 3.7, 0.28, "相关视频入口（点击卡片跳转）", size=16, color=COLORS["dark"], bold=True)
    add_video_card(
        slide,
        8.9,
        1.02,
        3.95,
        1.55,
        ASSETS / "vertical_farm.jpg",
        "垂直农业自动化示例",
        "Jungheinrich / YouTube",
        "https://www.youtube.com/watch?v=xq_S_5BkUCI",
    )
    add_video_card(
        slide,
        8.9,
        2.86,
        3.95,
        1.55,
        ASSETS / "workers_tablet.jpg",
        "PhenoBot 温室表型平台",
        "Wageningen University & Research / YouTube",
        "https://www.youtube.com/watch?v=Se6sROZoTUs",
    )
    add_video_card(
        slide,
        8.9,
        4.70,
        3.95,
        1.55,
        ASSETS / "conference_stage.jpg",
        "农业数字孪生演示",
        "University of Southern Denmark / YouTube",
        "https://www.youtube.com/watch?v=EyBqsKfwGCc",
    )
    add_text(
        slide,
        8.92,
        6.62,
        3.9,
        0.34,
        "图片：Pexels 实拍素材｜视频：Jungheinrich、WUR、SDU 官方公开链接",
        size=8.9,
        color=COLORS["muted"],
        font_name=FONT_EN,
    )
    add_page_mark(slide, "05/05")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    build()
